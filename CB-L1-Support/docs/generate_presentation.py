"""
generate_presentation.py
─────────────────────────
Generates a designed, 20-slide PowerPoint deck (CB_L1_Support_Chatbot_Overview.pptx)
for a management presentation on the CB L1 Support Chatbot project.

Design language: dark indigo/navy + white cards on a soft light-gray canvas,
icon badges, flow diagrams with connector arrows, comparison cards, a weighted
stat bar, and code blocks — not plain bullet lists / plain tables everywhere.

Run:  python docs/generate_presentation.py
Output: docs/CB_L1_Support_Chatbot_Overview.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

TOTAL = 20

# ── Palette ──────────────────────────────────────────────────────────────
NAVY = RGBColor(0x0B, 0x0F, 0x1E)
NAVY2 = RGBColor(0x15, 0x1B, 0x30)
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_DK = RGBColor(0x37, 0x2F, 0xA8)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
INK = RGBColor(0x11, 0x18, 0x27)
SLATE = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
BG = RGBColor(0xF4, 0xF6, 0xFB)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
SHADOW = RGBColor(0xDD, 0xE2, 0xEC)
RULE = RGBColor(0xE3, 0xE8, 0xF0)
CODE_BG = RGBColor(0x11, 0x18, 0x27)
CODE_TX = RGBColor(0xA5, 0xF3, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MX = Inches(0.55)
FULL_W = Inches(13.333 - 1.1)

prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
BLANK = prs.slide_layouts[6]
PAGE = [0]


# ── Primitives ───────────────────────────────────────────────────────────
def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def solid(shp, color, line=False, line_color=None, line_w=Pt(0.75)):
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line and line_color:
        shp.line.color.rgb = line_color
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def rect(slide, x, y, w, h, color, **kw):
    return solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h), color, **kw)


def oval(slide, x, y, w, h, color, **kw):
    return solid(slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h), color, **kw)


def rrect(slide, x, y, w, h, color, radius=0.08, **kw):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    return solid(shp, color, **kw)


def chevron(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    try:
        shp.adjustments[0] = 0.55
    except Exception:
        pass
    return solid(shp, color)


def card(slide, x, y, w, h, radius=0.06, fill=CARD, shadow_offset=Pt(4)):
    """White rounded card with a soft drop-shadow illusion behind it."""
    rrect(slide, Emu(int(x) + int(shadow_offset)), Emu(int(y) + int(shadow_offset)), w, h, SHADOW, radius=radius)
    return rrect(slide, x, y, w, h, fill, radius=radius)


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paragraphs = [runs] if runs and isinstance(runs[0], tuple) else runs
    p0 = tf.paragraphs[0]
    for i, para in enumerate(paragraphs):
        p = p0 if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (t, size, color, bold, italic) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = "Segoe UI"
    return tb


def icon_badge(slide, cx, cy, d, glyph, bgcolor, fgcolor=WHITE, size=15):
    oval(slide, Emu(int(cx - int(d) / 2)), Emu(int(cy - int(d) / 2)), d, d, bgcolor)
    text(slide, Emu(int(cx - int(d))), Emu(int(cy - int(d) / 2)), Emu(int(int(d) * 2)), d,
         [(glyph, size, fgcolor, True, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def pill(slide, x, y, label, bgcolor=NAVY2, fgcolor=WHITE, size=10.5, pad=0.14):
    w = Inches(pad * 2 + 0.1 * len(label))
    h = Inches(0.34)
    rrect(slide, x, y, w, h, bgcolor, radius=0.5)
    text(slide, x, y, w, h, [(label, size, fgcolor, True, False)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return w


def code_block(slide, x, y, w, h, lines, size=12.5):
    rrect(slide, x, y, w, h, CODE_BG, radius=0.07)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(10)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.name = "Consolas"
        r.font.color.rgb = CODE_TX


SECTION_ICONS = {
    "Overview": ("≡", INDIGO),
    "Context": ("!", AMBER),
    "Architecture": ("⌂", INDIGO),
    "Knowledge": ("▤", CYAN),
    "Backend": ("{ }", INDIGO_DK),
    "Feature": ("⚡", AMBER),
    "Frontend": ("▣", CYAN),
    "Config": ("#", INDIGO),
    "Deploy": ("▲", GREEN),
    "Playbook": ("✎", INDIGO_DK),
    "Wrap": ("✓", GREEN),
}


def header(slide, section_key, kicker, title):
    PAGE[0] += 1
    glyph, color = SECTION_ICONS.get(section_key, ("●", INDIGO))
    rect(slide, 0, 0, SLIDE_W, Inches(0.09), color)
    icon_badge(slide, Emu(int(MX) + int(Inches(0.28))), Inches(0.62), Inches(0.56), glyph, color, size=17)
    text(slide, Emu(int(MX) + int(Inches(0.66))), Inches(0.34), Inches(9), Inches(0.32),
         [(kicker.upper(), 12.5, color, True, False)])
    text(slide, Emu(int(MX) + int(Inches(0.66))), Inches(0.6), Inches(10.5), Inches(0.6),
         [(title, 24, INK, True, False)])
    chip_w = Inches(0.85)
    rrect(slide, Inches(12.1), Inches(0.36), chip_w, Inches(0.34), NAVY2, radius=0.5)
    text(slide, Inches(12.1), Inches(0.36), chip_w, Inches(0.34),
         [(f"{PAGE[0]:02d} / {TOTAL}", 10, WHITE, True, False)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide):
    rect(slide, MX, Inches(7.16), FULL_W, Pt(0.75), RULE)
    text(slide, MX, Inches(7.22), Inches(8), Inches(0.3),
         [("CB L1 Support Chatbot", 9, MUTED, True, False), ("   ·   Technical Overview", 9, MUTED, False, False)])


def new_slide(section_key, kicker, title):
    s = prs.slides.add_slide(BLANK)
    bg(s, BG)
    header(s, section_key, kicker, title)
    footer(s)
    return s


def bullet_list(slide, x, y, w, items, size=14, gap=Pt(9)):
    """items: (text, level, color|None) — level 0 = dot, 1 = dash (indented)."""
    tb = slide.shapes.add_textbox(x, y, w, Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for t, lvl, color in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = gap if lvl == 0 else Pt(4)
        mark = "●" if lvl == 0 else "–"
        indent = "" if lvl == 0 else "     "
        r = p.add_run()
        r.text = f"{indent}{mark}  {t}"
        r.font.size = Pt(size if lvl == 0 else size - 1.5)
        r.font.bold = lvl == 0 and color is not None
        r.font.color.rgb = color or (INK if lvl == 0 else SLATE)
        r.font.name = "Segoe UI"


# ══════════════════════════════════════════════════════════════════════
# 1. TITLE
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, SLIDE_W, Inches(0.09), INDIGO)
oval(s, Inches(10.6), Inches(-1.4), Inches(4.2), Inches(4.2), NAVY2)
oval(s, Inches(11.6), Inches(-0.6), Inches(2.6), Inches(2.6), INDIGO_DK)
oval(s, Inches(12.25), Inches(0.15), Inches(1.1), Inches(1.1), CYAN)
icon_badge(s, Inches(1.35), Inches(1.55), Inches(0.9), "CB", INDIGO, size=20)
text(s, Inches(0.95), Inches(2.55), Inches(3.6), Inches(0.4), [("ENGINEERING OVERVIEW", 13.5, CYAN, True, False)])
text(s, Inches(0.9), Inches(2.95), Inches(11.6), Inches(1.3), [("CB L1 Support Chatbot", 44, WHITE, True, False)])
text(s, Inches(0.95), Inches(3.95), Inches(11), Inches(0.6),
     [("Architecture, features & file-by-file reference", 18, RGBColor(0xCB, 0xD5, 0xE1), False, False)])
px, py = Inches(0.95), Inches(4.75)
for label in ["FastAPI", "React + TypeScript", "Qdrant", "LLM-agnostic"]:
    w = pill(s, px, py, label, bgcolor=NAVY2, fgcolor=RGBColor(0xE2, 0xE8, 0xF0))
    px = Emu(int(px) + int(w) + int(Inches(0.18)))
text(s, Inches(0.95), Inches(6.85), Inches(8), Inches(0.35),
     [("Defect Intelligence  +  Documentation RAG Assistant", 11, MUTED, False, True)])

# ══════════════════════════════════════════════════════════════════════
# 2. AGENDA
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Overview", "Overview", "Agenda")
items = [
    "Problem & solution", "System architecture", "The two knowledge bases",
    "Hybrid retrieval scoring", "Backend file reference", "Feature deep dive",
    "Frontend file reference", "Config & deployment", "Developer playbook",
    "Limitations & next steps",
]
col_w, row_h, gx, gy = Inches(5.85), Inches(0.86), MX, Inches(1.55)
for i, label in enumerate(items):
    col, row = divmod(i, 5)
    x = Emu(int(gx) + col * (int(col_w) + int(Inches(0.3))))
    y = Emu(int(gy) + row * (int(row_h) + int(Inches(0.12))))
    card(s, x, y, col_w, row_h, radius=0.18)
    icon_badge(s, Emu(int(x) + int(Inches(0.5))), Emu(int(y) + int(row_h) // 2), Inches(0.5), str(i + 1), INDIGO, size=15)
    text(s, Emu(int(x) + int(Inches(0.95))), y, Emu(int(col_w) - int(Inches(1.1))), row_h,
         [(label, 14.5, INK, True, False)], anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════════════════════════════════════════
# 3. PROBLEM & SOLUTION
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Context", "Business Context", "Problem & Solution")
cw, ch, cy = Inches(5.95), Inches(5.15), Inches(1.6)
card(s, MX, cy, cw, ch, radius=0.05)
icon_badge(s, Emu(int(MX) + int(Inches(0.55))), Emu(int(cy) + int(Inches(0.55))), Inches(0.6), "!", RED, size=20)
text(s, Emu(int(MX) + int(Inches(1.0))), Emu(int(cy) + int(Inches(0.3))), Inches(4.6), Inches(0.5),
     [("The Problem", 19, INK, True, False)])
bullet_list(s, Emu(int(MX) + int(Inches(0.45))), Emu(int(cy) + int(Inches(1.25))), Inches(5.1), [
    ("L1 engineers re-investigate defects already solved before", 0, None),
    ("Historical Jira fix knowledge is scattered across thousands of tickets", 0, None),
    ("Product documentation lives separately from real defect history", 0, None),
], size=14, gap=Pt(14))

cx2 = Emu(int(MX) + int(cw) + int(Inches(0.35)))
card(s, cx2, cy, cw, ch, radius=0.05)
icon_badge(s, Emu(int(cx2) + int(Inches(0.55))), Emu(int(cy) + int(Inches(0.55))), Inches(0.6), "✓", GREEN, size=20)
text(s, Emu(int(cx2) + int(Inches(1.0))), Emu(int(cy) + int(Inches(0.3))), Inches(4.6), Inches(0.5),
     [("The Solution", 19, INK, True, False)])
bullet_list(s, Emu(int(cx2) + int(Inches(0.45))), Emu(int(cy) + int(Inches(1.25))), Inches(5.1), [
    ("One assistant, two grounded evidence sources", 0, None),
    ("Every answer cites real Jira keys / document sources", 0, None),
    ("Knowledge bases are LIVE — add a defect or PDF, searchable instantly", 0, None),
], size=14, gap=Pt(14))

# ══════════════════════════════════════════════════════════════════════
# 4. ARCHITECTURE — FLOW DIAGRAM
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Architecture", "Architecture", "System Architecture — Request Flow")
steps = [("Frontend", "React SPA", CYAN), ("FastAPI", "Intent routing", INDIGO),
         ("Retrieval", "Qdrant hybrid", INDIGO_DK), ("LLM Synth.", "Cited answer", AMBER),
         ("Streamed", "SSE to UI", GREEN)]
n = len(steps)
box_w, box_h, gap_w = Inches(2.05), Inches(1.35), Inches(0.28)
total_w = n * int(box_w) + (n - 1) * int(gap_w)
start_x = int(MX) + (int(FULL_W) - total_w) // 2
y0 = Inches(2.35)
x = start_x
for i, (title_, sub, color) in enumerate(steps):
    card(s, Emu(x), y0, box_w, box_h, radius=0.12)
    icon_badge(s, Emu(x + int(box_w) // 2), Emu(int(y0) + int(Inches(0.42))), Inches(0.5), str(i + 1), color, size=15)
    text(s, Emu(x), Emu(int(y0) + int(Inches(0.78))), box_w, Inches(0.3),
         [(title_, 13.5, INK, True, False)], align=PP_ALIGN.CENTER)
    text(s, Emu(x), Emu(int(y0) + int(Inches(1.06))), box_w, Inches(0.25),
         [(sub, 10.5, SLATE, False, False)], align=PP_ALIGN.CENTER)
    x += int(box_w)
    if i < n - 1:
        chevron(s, Emu(x), Emu(int(y0) + int(Inches(0.45))), Emu(int(gap_w) + int(Inches(0.02))), Inches(0.45), MUTED)
        x += int(gap_w)

bullet_list(s, MX, Inches(4.15), FULL_W, [
    ("Two Qdrant collections on one cluster — \"defect\" (Jira history) and \"chatbot\" (documentation)", 0, None),
    ("LLM providers are pluggable — GitHub Models / Azure OpenAI / GitHub Copilot API via one config switch", 0, None),
    ("Frontend and API are served by the SAME FastAPI process — single origin, zero CORS complexity", 0, None),
    ("Jira Cloud REST API v3 is the live source for on-demand defect ingestion", 0, None),
], size=14.5, gap=Pt(10))

# ══════════════════════════════════════════════════════════════════════
# 5. TWO QDRANT COLLECTIONS — COMPARISON CARDS
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Knowledge", "Knowledge Bases", "Two Qdrant Collections")
cw, ch, cy = Inches(5.95), Inches(5.15), Inches(1.6)


def kb_card(x, title_, icon, color, rows):
    card(s, x, cy, cw, ch, radius=0.05)
    icon_badge(s, Emu(int(x) + int(Inches(0.5))), Emu(int(cy) + int(Inches(0.5))), Inches(0.55), icon, color, size=16)
    text(s, Emu(int(x) + int(Inches(0.95))), Emu(int(cy) + int(Inches(0.28))), Inches(4.5), Inches(0.5),
         [(title_, 18, INK, True, False)])
    ry = Emu(int(cy) + int(Inches(1.15)))
    for label, val in rows:
        text(s, Emu(int(x) + int(Inches(0.4))), ry, Inches(1.6), Inches(0.42),
             [(label, 11, MUTED, True, False)])
        text(s, Emu(int(x) + int(Inches(2.0))), ry, Emu(int(cw) - int(Inches(2.3))), Inches(0.55),
             [(val, 12, INK, False, False)])
        ry = Emu(int(ry) + int(Inches(0.66)))


kb_card(MX, '"defect" collection', "▤", INDIGO, [
    ("Content", "Historical Jira defects — root cause, fix, comments"),
    ("Populated", "Offline build + on-demand Add Defect (key/JQL)"),
    ("Chunking", "5 types/defect (problem, cause, fix, comment, context)"),
    ("Retrieval", "Hybrid: semantic + full-text + metadata + quality"),
    ("Lookup", "Exact Jira key overrides search (e.g. HCBS-95506)"),
])
kb_card(Emu(int(MX) + int(cw) + int(Inches(0.35))), '"chatbot" collection', "▣", CYAN, [
    ("Content", "Product documentation — PDF text + tables"),
    ("Populated", "Offline build + on-demand PDF upload"),
    ("Chunking", "Section-aware text + table rows/fragments"),
    ("Retrieval", "Semantic top-K + ±1 neighbor table context"),
    ("Lookup", "N/A — no key-based equivalent"),
])

# ══════════════════════════════════════════════════════════════════════
# 6. HYBRID RETRIEVAL SCORING — WEIGHTED STAT BAR
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Knowledge", "Knowledge Bases", "Hybrid Defect Retrieval Scoring")
weights = [("Semantic", 0.45, INDIGO), ("Keyword", 0.25, CYAN), ("Metadata", 0.15, AMBER),
           ("Quality", 0.10, GREEN), ("Fixed-boost", 0.05, RED)]
bar_x, bar_y, bar_w, bar_h = MX, Inches(1.75), FULL_W, Inches(0.62)
x = int(bar_x)
for label, wgt, color in weights:
    seg_w = int(int(bar_w) * wgt)
    rect(s, Emu(x), bar_y, Emu(seg_w), bar_h, color)
    if wgt >= 0.10:
        text(s, Emu(x), bar_y, Emu(seg_w), bar_h, [(f"{int(wgt*100)}%", 15, WHITE, True, False)],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += seg_w
lx = int(bar_x)
ly = Inches(2.55)
for label, wgt, color in weights:
    oval(s, Emu(lx), ly, Inches(0.18), Inches(0.18), color)
    text(s, Emu(lx + int(Inches(0.26))), Emu(int(ly) - int(Inches(0.03))), Inches(1.7), Inches(0.3),
         [(f"{label} {int(wgt*100)}%", 11.5, INK, True, False)])
    lx += int(Inches(2.35))

bullet_list(s, MX, Inches(3.15), FULL_W, [
    ("Semantic — Qdrant cosine similarity on the embedding vector, top-60 candidates", 0, None),
    ("Keyword — Qdrant full-text index + in-memory BM25-style re-score for exact-term matches", 0, None),
    ("Metadata — bonus for matching components, labels, fix pattern, team", 0, None),
    ("Quality — bonus for richer, higher-confidence extracted defects", 0, None),
    ("Fixed-boost — promotes confirmed-fixed defects; demotes cancelled ones", 0, None),
], size=13.5, gap=Pt(7))

card(s, MX, Inches(5.85), FULL_W, Inches(0.9), radius=0.1, fill=RGBColor(0xFF, 0xF7, 0xE6))
icon_badge(s, Emu(int(MX) + int(Inches(0.5))), Emu(int(Inches(5.85)) + int(Inches(0.45))), Inches(0.44), "!", AMBER, size=14)
text(s, Emu(int(MX) + int(Inches(0.85))), Inches(5.98), Emu(int(FULL_W) - int(Inches(1.1))), Inches(0.65),
     [("Semantic-only search caused a measurable precision drop in testing — the keyword signal (Qdrant full-text) restored accuracy while staying fully live-updatable.", 12, INK, False, False)])

# ══════════════════════════════════════════════════════════════════════
# 7-10. BACKEND FILE TABLES (styled cards, not plain pptx tables)
# ══════════════════════════════════════════════════════════════════════
def file_rows(s, x, y, w, rows, row_h=Inches(0.62), name_w=Inches(2.9), size=11.5):
    ry = y
    for i, (name, desc) in enumerate(rows):
        fill = CARD if i % 2 == 0 else RGBColor(0xF7, 0xF9, 0xFD)
        rrect(s, x, ry, w, row_h, fill, radius=0.12, line=True, line_color=RULE, line_w=Pt(0.75))
        text(s, Emu(int(x) + int(Inches(0.22))), ry, Emu(int(name_w) - int(Inches(0.2))), row_h,
             [(name, size, INDIGO_DK, True, False)], anchor=MSO_ANCHOR.MIDDLE)
        text(s, Emu(int(x) + int(name_w) + int(Inches(0.1))), ry, Emu(int(w) - int(name_w) - int(Inches(0.4))), row_h,
             [(desc, size, SLATE, False, False)], anchor=MSO_ANCHOR.MIDDLE)
        ry = Emu(int(ry) + int(row_h) + int(Pt(4)))


s = new_slide("Backend", "Backend", "Core Config, LLM Provider & API Layer")
file_rows(s, MX, Inches(1.6), FULL_W, [
    ("config.py", "Single source of truth for ALL env vars: Jira, Qdrant, LLM provider, retrieval tuning, feature flags"),
    ("llm_provider.py", "Provider-agnostic LLM abstraction — switches Azure / GitHub Models / Copilot via LLM_PROVIDER"),
    ("run_chatbot.py", "CLI entry point: --build (offline KB build), --serve (start API), --ask (one-off query)"),
    ("api/app.py", "FastAPI app factory — mounts built SPA, includes routes, pre-warms Qdrant/LLM on startup"),
    ("api/routes.py", "All HTTP endpoints: /health, /api/ask(/stream), /api/rag/ingest(/stream), /api/defects/add(-jql)"),
])

s = new_slide("Backend", "Backend", "Chatbot Core Pipeline (chatbot/)")
file_rows(s, MX, Inches(1.55), FULL_W, [
    ("schemas.py", "DefectRecord dataclass + shared constants (fix patterns, chunk types)"),
    ("intent_router.py", "Rule-based classifier: defect-by-key / diagnostic / similar-defects / general help"),
    ("prompts.py", "LLM prompt templates enforcing evidence-grounded, cited, structured answers"),
    ("answer_generator.py", "Synthesizes retrieved evidence into an answer via LLM, with deterministic fallback"),
    ("retriever.py", "Hybrid retriever — Qdrant semantic + keyword + metadata + quality scoring"),
    ("defect_qa.py", "Top-level orchestrator: routes intent, retrieves evidence, calls doc-RAG, streams answer"),
    ("utils.py", "Shared helpers — JSON I/O, Jira field shims, PII masking, key extraction, tokenization"),
], row_h=Inches(0.56), size=11)

s = new_slide("Backend", "Backend", "Knowledge-Base Build & On-Demand Ingestion")
file_rows(s, MX, Inches(1.55), FULL_W, [
    ("build_knowledge_base.py", "Normalizes raw Jira data, extracts root cause / fix / pattern / type, masks PII"),
    ("build_chunks.py", "Splits each defect into 5 typed chunks for embedding"),
    ("build_embeddings.py", "Embeds all chunks (remote provider or local fallback), L2-normalizes vectors"),
    ("migrate_defects_qdrant.py", "One-time / re-sync migration of embeddings into the Qdrant \"defect\" collection"),
    ("defect_ingest.py", "On-demand: fetches from Jira live (single key or bulk JQL), runs the SAME pipeline as offline build"),
    ("defect_vector_store.py", "Qdrant wrapper — semantic + full-text keyword search, idempotent upsert, delete-by-key"),
], row_h=Inches(0.63), size=11)

s = new_slide("Backend", "Backend", "Documentation RAG Pipeline (rag_docs/)")
file_rows(s, MX, Inches(1.55), FULL_W, [
    ("document_ingestor.py", "Parses PDFs via LlamaParse into structured TEXT/TABLE elements with section context"),
    ("semantic_chunker.py / table_parser.py", "Splits parsed elements into semantic chunks (text + table rows/fragments)"),
    ("vector_store.py", "Qdrant wrapper for the \"chatbot\" collection — embed, upsert, search, ±1 context expansion"),
    ("answer_generator.py", "Formats retrieved chunks and calls the LLM for a cited answer ([Source: file.pdf, Page: X])"),
    ("query_processor.py / query_router.py / rag_agent.py", "Query enhancement, context routing, and pipeline coordination"),
    ("llm_factory.py / pipeline.py / ingest.py / models.py", "RAG-scoped provider selection, cached pipeline builder, PDF ingestion entrypoint, shared models"),
], row_h=Inches(0.63), size=10.5)

# ══════════════════════════════════════════════════════════════════════
# 11. FEATURE DEEP DIVE — TWO FLOW PIPELINES
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Feature", "Feature Deep Dive", "Add Defect  &  Upload PDF")


def pipeline_row(s, y, label, icon, color, steps):
    icon_badge(s, Emu(int(MX) + int(Inches(0.3))), Emu(int(y) + int(Inches(0.42))), Inches(0.6), icon, color, size=18)
    text(s, Emu(int(MX) + int(Inches(0.72))), y, Inches(2.3), Inches(0.85),
         [(label, 15, INK, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    n = len(steps)
    box_w, box_h, gap_w = Inches(1.55), Inches(0.85), Inches(0.22)
    x = int(MX) + int(Inches(2.75))
    for i, st in enumerate(steps):
        card(s, Emu(x), y, box_w, box_h, radius=0.16)
        text(s, Emu(x), y, box_w, box_h, [(st, 11.5, INK, True, False)],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += int(box_w)
        if i < n - 1:
            chevron(s, Emu(x), Emu(int(y) + int(Inches(0.28))), Emu(int(gap_w)), Inches(0.3), MUTED)
            x += int(gap_w)


pipeline_row(s, Inches(1.75), "Add Defect", "+", INDIGO,
             ["Fetch (Jira)", "Normalize", "Chunk (×5)", "Embed", "Upsert"])
pipeline_row(s, Inches(3.05), "Upload PDF", "⇪", CYAN,
             ["Parse", "Chunk", "Embed", "Upsert"])

card(s, MX, Inches(4.55), FULL_W, Inches(2.15), radius=0.05)
bullet_list(s, Emu(int(MX) + int(Inches(0.35))), Inches(4.8), Emu(int(FULL_W) - int(Inches(0.7))), [
    ("Single Jira key OR bulk JQL query (capped at 200 issues/run) — same normalization pipeline as offline build", 0, None),
    ("Live SSE progress streamed to the UI at every stage", 0, None),
    ("PDF upload keeps section headings + table structure; re-uploading the same file overwrites, never duplicates", 0, None),
    ("Both features update the in-memory retriever instantly — searchable on the very next question, no restart", 0, INDIGO_DK),
], size=13.5, gap=Pt(10))

# ══════════════════════════════════════════════════════════════════════
# 12-14. FRONTEND FILE TABLES
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Frontend", "Frontend", "App Shell, Routing & Pages")
file_rows(s, MX, Inches(1.6), FULL_W, [
    ("main.tsx / App.tsx", "Entry point; wraps app in QueryClient / Theme / Router / Chat providers"),
    ("routes/AppRoutes.tsx", "Lazy routes: \"/\" → Landing, \"/chat\" → Chat, animated transitions"),
    ("layouts/ChatLayout.tsx", "Responsive shell — collapsible sidebar (desktop) / drawer (mobile) + header"),
    ("pages/ChatPage.tsx", "Wraps ChatLayout + ChatWindow — the main chat screen"),
    ("pages/LandingPage.tsx + landing/*", "Public marketing page (Hero, Features, HowItWorks, Footer, etc.)"),
])

s = new_slide("Frontend", "Frontend", "Chat Components (components/chat/)")
file_rows(s, MX, Inches(1.55), FULL_W, [
    ("ChatWindow / MessageList / MessageBubble", "Message container, scrollable list, individual bubble with copy/regenerate"),
    ("ChatInput.tsx", "Auto-resizing input, 2000-char limit, send/stop, keyboard shortcuts"),
    ("OptionsCard.tsx", "Lets user pick answer sections (root cause / resolve / similar)"),
    ("SourcesPanel.tsx", "Collapsible panel of similar defects with status badges & relevance scores"),
    ("Markdown.tsx", "Renders answer markdown with code/table formatting"),
    ("TypingIndicator.tsx / EmptyState.tsx", "Streaming animation; first-run welcome screen with example prompts"),
], row_h=Inches(0.58), size=11)

s = new_slide("Frontend", "Frontend", "Header, Global State, Services & Types")
file_rows(s, MX, Inches(1.55), FULL_W, [
    ("components/layout/Header.tsx", "Connection badge, theme toggle, Upload PDF button, Add Defect popover (key/JQL modes)"),
    ("components/layout/Sidebar.tsx", "Conversation list — new/delete/pin/search, persisted to localStorage"),
    ("features/chat/ChatProvider.tsx", "Global chat state — conversations, streaming, send/regenerate/stop"),
    ("services/api.ts", "HTTP client for every backend endpoint incl. SSE streaming"),
    ("types/index.ts", "TypeScript contracts mirroring the FastAPI backend exactly"),
    ("hooks/, contexts/, utils/", "Health polling, clipboard, responsive breakpoints, theme, formatting, storage"),
], row_h=Inches(0.58), size=11)

# ══════════════════════════════════════════════════════════════════════
# 15. ENV VARS
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Config", "Configuration", "Environment Variables & Secrets")
file_rows(s, MX, Inches(1.55), FULL_W, [
    ("Qdrant (required)", "QDRANT_URL, QDRANT_API_KEY, DEFECT_QDRANT_COLLECTION, QDRANT_COLLECTION, DEFECT_VECTOR_DIM"),
    ("Jira (required for ingestion)", "JIRA_URL, JIRA_USERNAME, JIRA_API_TOKEN, JIRA_BASE_URL"),
    ("LLM provider switch", "LLM_PROVIDER, EMBED_PROVIDER, RAG_LLM_PROVIDER, RAG_EMBED_PROVIDER"),
    ("Provider credentials", "AZURE_OPENAI_*  /  GITHUB_TOKEN  /  COPILOT_TOKEN, COPILOT_API_ENDPOINT"),
    ("LlamaParse", "LLAMA_CLOUD_API_KEY (PDF ingestion)"),
    ("Retrieval tuning & flags", "TOP_K_RESULTS, RELEVANCE_MIN_SCORE, USE_LLM, USE_RAG_DOCS, MASK_SENSITIVE_DATA"),
    ("Ingestion limits", "JQL_ADD_MAX_RESULTS (200), RAG_INGEST_MAX_MB (20)"),
], row_h=Inches(0.62), size=11.5)

# ══════════════════════════════════════════════════════════════════════
# 16. DEPLOYMENT — TWO CARDS + CODE BLOCK
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Deploy", "Deployment", "Packaging & Deployment")
cw, ch, cy = Inches(5.95), Inches(2.35), Inches(1.55)
card(s, MX, cy, cw, ch, radius=0.06)
icon_badge(s, Emu(int(MX) + int(Inches(0.5))), Emu(int(cy) + int(Inches(0.45))), Inches(0.5), "▣", INDIGO, size=15)
text(s, Emu(int(MX) + int(Inches(0.9))), Emu(int(cy) + int(Inches(0.22))), Inches(4.5), Inches(0.45),
     [("Docker (multi-stage)", 15.5, INK, True, False)])
bullet_list(s, Emu(int(MX) + int(Inches(0.4))), Emu(int(cy) + int(Inches(0.95))), Emu(int(cw) - int(Inches(0.7))), [
    ("Stage 1 (Node 20) builds the React SPA", 0, None),
    ("Stage 2 (Python 3.12-slim) serves API + SPA on one port", 0, None),
], size=12, gap=Pt(6))

cx2 = Emu(int(MX) + int(cw) + int(Inches(0.35)))
card(s, cx2, cy, cw, ch, radius=0.06)
icon_badge(s, Emu(int(cx2) + int(Inches(0.5))), Emu(int(cy) + int(Inches(0.45))), Inches(0.5), "☁", GREEN, size=15)
text(s, Emu(int(cx2) + int(Inches(0.9))), Emu(int(cy) + int(Inches(0.22))), Inches(4.5), Inches(0.45),
     [("Render (native Python)", 15.5, INK, True, False)])
bullet_list(s, Emu(int(cx2) + int(Inches(0.4))), Emu(int(cy) + int(Inches(0.95))), Emu(int(cw) - int(Inches(0.7))), [
    ("Free tier, no payment card required", 0, None),
    ("Frontend pre-built locally & committed to git", 0, None),
], size=12, gap=Pt(6))

code_block(s, MX, Inches(4.15), FULL_W, Inches(1.55), [
    "$ docker build -t cb-l1-support-chatbot:latest .",
    "$ docker run -d -p 5100:10000 -v <path>.env:/app/.env:ro cb-l1-support-chatbot:latest",
])
text(s, MX, Inches(5.9), FULL_W, Inches(0.5),
     [("Verified locally — health check confirms Qdrant + LLM + RAG reachable inside the container.", 12.5, GREEN, True, True)])

# ══════════════════════════════════════════════════════════════════════
# 17. DEVELOPER PLAYBOOK
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Playbook", "Developer Playbook", "\"I Want to Change X\" → Edit This File")
file_rows(s, MX, Inches(1.55), FULL_W, [
    ("Change similar-defect ranking weights", "backend/chatbot/retriever.py"),
    ("Change the LLM's answer style/structure", "backend/chatbot/prompts.py"),
    ("Switch LLM provider (Azure ↔ GitHub ↔ Copilot)", "backend/.env → LLM_PROVIDER"),
    ("Change how questions are classified", "backend/chatbot/intent_router.py"),
    ("Add a new API endpoint", "backend/api/routes.py"),
    ("Change the Add Defect popover / UI text", "frontend/src/components/layout/Header.tsx"),
    ("After ANY frontend change", "cd frontend && npm run build"),
], row_h=Inches(0.56), name_w=Inches(4.6), size=11.5)

# ══════════════════════════════════════════════════════════════════════
# 18. LIMITATIONS & NEXT STEPS
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Wrap", "Wrap-Up", "Limitations & Next Steps")
cw, ch, cy = Inches(5.95), Inches(5.15), Inches(1.6)
card(s, MX, cy, cw, ch, radius=0.05, fill=RGBColor(0xFF, 0xF7, 0xE6))
icon_badge(s, Emu(int(MX) + int(Inches(0.5))), Emu(int(cy) + int(Inches(0.5))), Inches(0.55), "!", AMBER, size=17)
text(s, Emu(int(MX) + int(Inches(0.95))), Emu(int(cy) + int(Inches(0.28))), Inches(4.5), Inches(0.5),
     [("Known Limitations", 17, INK, True, False)])
bullet_list(s, Emu(int(MX) + int(Inches(0.4))), Emu(int(cy) + int(Inches(1.1))), Emu(int(cw) - int(Inches(0.7))), [
    ("Chat history is client-side only (localStorage) — no cross-device continuity", 0, None),
    ("Local defect_knowledge_base.json needs a persistent volume in Docker to survive restarts", 0, None),
    ("Jira API token must be verified/rotated periodically", 0, None),
    ("JQL bulk-add is capped at 200 issues/run as a safety guard", 0, None),
], size=13, gap=Pt(11))

cx2 = Emu(int(MX) + int(cw) + int(Inches(0.35)))
card(s, cx2, cy, cw, ch, radius=0.05, fill=RGBColor(0xEC, 0xFD, 0xF5))
icon_badge(s, Emu(int(cx2) + int(Inches(0.5))), Emu(int(cy) + int(Inches(0.5))), Inches(0.55), "✓", GREEN, size=17)
text(s, Emu(int(cx2) + int(Inches(0.95))), Emu(int(cy) + int(Inches(0.28))), Inches(4.5), Inches(0.5),
     [("Recommended Next Steps", 17, INK, True, False)])
bullet_list(s, Emu(int(cx2) + int(Inches(0.4))), Emu(int(cy) + int(Inches(1.1))), Emu(int(cw) - int(Inches(0.7))), [
    ("Mount a persistent volume for backend/data/chatbot/ in production", 0, None),
    ("Verify/rotate Jira API credentials before relying on live ingestion", 0, None),
    ("Add regression tests around the retrieval scoring formula", 0, None),
    ("Periodically re-run the offline batch build as Jira terminology evolves", 0, None),
], size=13, gap=Pt(11))

# ══════════════════════════════════════════════════════════════════════
# 19. KEY TAKEAWAYS
# ══════════════════════════════════════════════════════════════════════
s = new_slide("Wrap", "Wrap-Up", "Key Takeaways")
takeaways = [
    ("★", INDIGO, "A single assistant unifies historical defect knowledge and product documentation"),
    ("≈", CYAN, "Retrieval is hybrid — semantic + keyword + metadata — tuned to preserve accuracy, not just \"vector search\""),
    ("⚡", AMBER, "Both knowledge bases are LIVE-updatable from the UI — no rebuilds, no downtime"),
    ("⌘", GREEN, "Provider-agnostic for LLMs; deployable via Docker or Render with minimal changes"),
]
ty = Inches(1.75)
for glyph, color, txt_ in takeaways:
    card(s, MX, ty, FULL_W, Inches(1.0), radius=0.12)
    icon_badge(s, Emu(int(MX) + int(Inches(0.55))), Emu(int(ty) + int(Inches(0.5))), Inches(0.62), glyph, color, size=18)
    text(s, Emu(int(MX) + int(Inches(1.1))), ty, Emu(int(FULL_W) - int(Inches(1.4))), Inches(1.0),
         [(txt_, 15, INK, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    ty = Emu(int(ty) + int(Inches(1.15)))

card(s, MX, ty, FULL_W, Inches(0.85), radius=0.12, fill=NAVY2)
text(s, Emu(int(MX) + int(Inches(0.4))), ty, Emu(int(FULL_W) - int(Inches(0.8))), Inches(0.85),
     [("This deck + the developer playbook is the reference for making future changes safely.", 14.5, WHITE, True, True)],
     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# 20. CLOSING
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, SLIDE_W, Inches(0.09), INDIGO)
oval(s, Inches(-1.2), Inches(4.6), Inches(4), Inches(4), NAVY2)
oval(s, Inches(-0.3), Inches(5.3), Inches(2.4), Inches(2.4), INDIGO_DK)
icon_badge(s, Inches(1.35), Inches(2.55), Inches(0.9), "CB", INDIGO, size=20)
text(s, Inches(0.9), Inches(3.35), Inches(11), Inches(1), [("Thank You", 42, WHITE, True, False)])
text(s, Inches(0.95), Inches(4.2), Inches(11), Inches(0.6), [("Questions & Discussion", 18, CYAN, False, False)])
text(s, Inches(0.95), Inches(6.9), Inches(11), Inches(0.35),
     [("CB L1 Support Chatbot  ·  Technical Overview", 10.5, MUTED, False, True)])

out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CB_L1_Support_Chatbot_Overview.pptx")
try:
    prs.save(out_path)
except PermissionError:
    out_path = out_path.replace(".pptx", "_new.pptx")
    prs.save(out_path)
    print("NOTE: original file was locked — saved as _new.pptx instead.")
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
